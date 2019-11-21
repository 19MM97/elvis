# -*- coding: utf-8 -*-
"""
Created on Sat Apr 27 19:06:23 2019

@author: draz
"""
import datetime as dt
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData, XyChartData
from pptx.chart.chart import Chart 
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import pandas as pd 
# import matplotlib.pyplot as plt
import calendar


class GetSlides:
    def __init__(self):
        self.prs = Presentation()
        self.get_cover_slide()
        
    def get_table(self, data, heading): 
        keys = [key for key in data.keys()]
        values = [value for value in data.values()]
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = heading
        x, y, cx, cy = Inches(0.5), Inches(3), Inches(9), Inches(2)
        shape = slide.shapes.add_table(2, len(keys), x, y, cx, cy)
        for c in range(len(keys)): 
            shape.table.cell(0, c).text = keys[c]
            shape.table.cell(1, c).text = str(values[c])

    def get_chart(self, data, heading): 
        # chart
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = heading
        chart_data = ChartData()
        chart_data.categories = data.index
        for series in data.columns: 
            chart_data.add_series(series, data[series].values)
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, 
            x, y, cx, cy, 
            chart_data
            ).chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.series[0].smooth = True 

    def get_chart_stack(self, heading, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = heading
        chart_data = CategoryChartData()
        chart_data.categories = data.keys()
        chart_data.add_series('served by each LP', data.values())
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, 
            x, 
            y, 
            cx, 
            cy, 
            chart_data,
            )
        chart.has_legend = True

    def get_cover_slide(self):
        # slide title
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = 'FlexNet4E-Mobility JF'
        # subtitle1.text = 'Mahmoud Draz, Marcus Voß'
        subtitle.text = 'Datum: %s' % dt.datetime.today().strftime('%d.%m.%Y')
        
    def get_pic_slide(self, heading, name): 
        # slide 1 picture
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = heading
        left = Inches(0.5)
        top = Inches(2)
        height = Inches(5.5)
        width = Inches(9)
        slide.shapes.add_picture(name, left, top, height=height, width=width)
        
    def get_graphics(self, control, p, n):
        df_load = pd.read_csv('LoadProfile_%s_%s_%s.csv' % (control, p, n))
        df_load['Time'] = pd.to_datetime(df_load['Time'])
        
        df_load['week_of_year'] = df_load['Time'].dt.week
        df_load['day_of_week'] = df_load['Time'].dt.day_name()
        df_load['quarter-of-hour'] = df_load['Time'].dt.quarter
        df_load['hour_of_day'] = df_load['Time'].dt.hour
        df_load['weekday'] = df_load['Time'].dt.weekday
        df_load['workingday/weekend'] = df_load['weekday']//5 
        
        hours = [dt.time(i).strftime('%H:%M') for i in range(24)]
#        days = [""]+list(calendar.day_abbr)
        
        ax = df_load['Total_Trafo_LP'][:10080].plot(title='week load profile')
        ax.set_xticklabels(['Mon', 'Tue', 'Wed', 'Thu', 'Fri', 'Sat', 'Sun'], rotation=0)
        ax.set_ylabel('Power/kVA')
        ax.set_xlabel('Day')
        plt = ax.get_figure()
        plt.savefig('week_LP_%s.png' % control, bbox_inches='tight')
        
        ax = df_load.groupby(['day_of_week', 'hour_of_day']).Total_Trafo_LP.mean().\
            unstack().plot(legend=False, kind='box', title='load ranges over the day')
        ax.set_xticks(range(0, 25, 1))
        ax.set_xticklabels(hours, rotation=50)
        ax.set_ylabel('Power/kW')
        ax.set_xlabel('Time')
        plt = ax.get_figure()
        plt.savefig('LP_agg_byhour_we_%s.png' % control, bbox_inches='tight')
        
        ax = df_load.groupby(['workingday/weekend', 'hour_of_day']).Total_Trafo_LP.mean().\
            unstack().plot(legend=False, kind='box', title='load ranges over the day')
        ax.set_xticks(range(0, 25, 1))
        ax.set_xticklabels(hours, rotation=50)
        ax.set_ylabel('Power/kW')
        ax.set_xlabel('Time')
        plt = ax.get_figure()
        plt.savefig('LP_agg_byhour_%s.png' % control, bbox_inches='tight')
        
        ax = df_load.groupby(['week_of_year', 'day_of_week']).Total_Trafo_LP.mean().\
            unstack().plot(legend=False, kind='box', title='load ranges over the days of the week')
        ax.set_xticks(range(1, 8, 1))
        ax.set_ylabel('Power/kW')
        ax.set_xlabel('Day')
        ax.set_xticklabels(['Thu', 'Fri', 'Sat', 'Sun', 'Mon', 'Tue', 'Wed'], rotation=0)
        plt = ax.get_figure()
        plt.savefig('LP_agg_byday_%s.png' % control, bbox_inches='tight')
        
        ax = df_load.groupby(['day_of_week', 'week_of_year']).Total_Trafo_LP.mean().\
            unstack().plot(legend=False, kind='box', title='load ranges the weeks of the year')
        ax.set_xticks(range(0, 53, 5))
        ax.set_ylabel('Power/kW')
        ax.set_xlabel('Week')
        ax.set_xticklabels(['W_%s' % n for n in range(0, 52, 5)], rotation=0)
        plt = ax.get_figure()
        plt.savefig('LP_agg_byweek_%s.png' % control, bbox_inches='tight')
        
    def get_slides(self, assumptions, indicators, control, p, n):
        self.get_graphics(control, p, n)
        self.get_table(assumptions, 'Szenario %s_%s: Annahmen' % (assumptions['Ort'], control))
        self.get_pic_slide('Szenario %s_%s: selected week load profile' %
                           (assumptions['Ort'], control), 'week_LP_%s.png' % control)
        self.get_pic_slide('Szenario %s_%s: year profile grouped by hour of the day' %
                           (assumptions['Ort'], control), 'LP_agg_byhour_%s.png' % control)
        self.get_pic_slide('Szenario %s_%s: year profile grouped by hour of the day_WE/WD' %
                           (assumptions['Ort'], control), 'LP_agg_byhour_we_%s.png' % control)
        self.get_pic_slide('Szenario %s_%s: year profile grouped by day of the week' %
                           (assumptions['Ort'], control), 'LP_agg_byday_%s.png' % control)
        self.get_pic_slide('Szenario %s_%s: year profile grouped byweek of the year' %
                           (assumptions['Ort'], control), 'LP_agg_byweek_%s.png' % control)
        self.get_table(heading='Szenario %s_%s: Kennzahlen' % (assumptions['Ort'], control), data=indicators)
        self.prs.save('Flexnet_JF_%s.pptx' % dt.datetime.today().strftime('%d%m%Y'))
        
# df_occupancy.to_csv('df_occupancy%s.csv'%control)
# plot1 = evLoad.plot(legend = False, title = 'EV Lastprofil')
# plot1.set(xlabel='Zeit', ylabel='Last (kWA)')
# fig1= plot1.get_figure()
# fig1.savefig('evlast_%s.png'%control,bbox_inches='tight')
#
# self.get_chart(evLoad,'Szenario %s_%s: EV Lastprofil'%(assumptions['Ort'], control))
# self.get_pic_slide('Szenario %s_%s: EV Lastprofil'%(assumptions['Ort'],control),'evlast_%s.png'%control)
# self.get_chart_stack(heading='Szenario %s_%s: Wie viele Autos pro LP? ' %
# (assumptions['Ort'],control), data=servedByEach)
